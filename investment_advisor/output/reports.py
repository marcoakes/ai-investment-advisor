"""
Report generation tools for creating PDF and PowerPoint outputs.
"""

from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import pandas as pd
from typing import Dict, Any, Optional, List
import os
from datetime import datetime
from ..core.base import BaseTool, ToolResult, ToolType


class PDFReportGenerator(BaseTool):
    """Generate comprehensive PDF reports."""
    
    def __init__(self, output_dir: str = "reports"):
        super().__init__("pdf_report_generator", ToolType.OUTPUT)
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
    
    def _setup_custom_styles(self):
        """Set up custom paragraph styles."""
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            textColor=colors.darkblue,
            alignment=TA_CENTER
        ))
        
        self.styles.add(ParagraphStyle(
            name='SectionHeader',
            parent=self.styles['Heading2'],
            fontSize=16,
            spaceAfter=12,
            textColor=colors.darkgreen,
            borderWidth=1,
            borderColor=colors.darkgreen,
            borderPadding=5
        ))
    
    def execute(self, report_data: Dict[str, Any], report_title: str = "Investment Analysis Report", 
                save_path: str = None, **kwargs) -> ToolResult:
        """Generate a PDF report."""
        try:
            if save_path is None:
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                save_path = os.path.join(self.output_dir, f"investment_report_{timestamp}.pdf")
            
            doc = SimpleDocTemplate(save_path, pagesize=A4, topMargin=1*inch)
            story = []
            
            # Title
            story.append(Paragraph(report_title, self.styles['CustomTitle']))
            story.append(Spacer(1, 20))
            
            # Executive Summary
            if 'executive_summary' in report_data:
                story.append(Paragraph("Executive Summary", self.styles['SectionHeader']))
                story.append(Paragraph(report_data['executive_summary'], self.styles['Normal']))
                story.append(Spacer(1, 12))
            
            # Stock Information
            if 'stock_info' in report_data:
                story = self._add_stock_info_section(story, report_data['stock_info'])
            
            # Performance Metrics
            if 'performance_metrics' in report_data:
                story = self._add_performance_section(story, report_data['performance_metrics'])
            
            # Technical Analysis
            if 'technical_analysis' in report_data:
                story = self._add_technical_section(story, report_data['technical_analysis'])
            
            # Charts
            if 'charts' in report_data:
                story = self._add_charts_section(story, report_data['charts'])
            
            # Risk Analysis
            if 'risk_analysis' in report_data:
                story = self._add_risk_section(story, report_data['risk_analysis'])
            
            # Recommendations
            if 'recommendations' in report_data:
                story = self._add_recommendations_section(story, report_data['recommendations'])
            
            # Build PDF
            doc.build(story)
            
            return ToolResult(
                success=True,
                data={'report_path': save_path},
                metadata={'report_type': 'pdf', 'title': report_title}
            )
        
        except Exception as e:
            return ToolResult(
                success=False,
                error=f"PDF report generation failed: {str(e)}"
            )
    
    def _add_stock_info_section(self, story: List, stock_info: Dict[str, Any]) -> List:
        """Add stock information section."""
        story.append(Paragraph("Stock Information", self.styles['SectionHeader']))
        
        info_data = [
            ['Symbol:', stock_info.get('symbol', 'N/A')],
            ['Company Name:', stock_info.get('company_name', 'N/A')],
            ['Sector:', stock_info.get('sector', 'N/A')],
            ['Industry:', stock_info.get('industry', 'N/A')],
            ['Market Cap:', stock_info.get('market_cap', 'N/A')],
            ['Current Price:', f"${stock_info.get('current_price', 0):.2f}"]
        ]
        
        table = Table(info_data, colWidths=[2*inch, 3*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        story.append(table)
        story.append(Spacer(1, 12))
        return story
    
    def _add_performance_section(self, story: List, performance: Dict[str, Any]) -> List:
        """Add performance metrics section."""
        story.append(Paragraph("Performance Metrics", self.styles['SectionHeader']))
        
        perf_data = [
            ['Metric', 'Value'],
            ['Total Return', f"{performance.get('total_return', 0):.2f}%"],
            ['Annualized Return', f"{performance.get('annualized_return', 0):.2f}%"],
            ['Sharpe Ratio', f"{performance.get('sharpe_ratio', 0):.3f}"],
            ['Maximum Drawdown', f"{performance.get('max_drawdown', 0):.2f}%"],
            ['Volatility', f"{performance.get('volatility', 0):.2f}%"],
            ['Win Rate', f"{performance.get('win_rate', 0):.2f}%"]
        ]
        
        table = Table(perf_data, colWidths=[2.5*inch, 2.5*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.darkblue),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.lightblue),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        story.append(table)
        story.append(Spacer(1, 12))
        return story
    
    def _add_technical_section(self, story: List, technical: Dict[str, Any]) -> List:
        """Add technical analysis section."""
        story.append(Paragraph("Technical Analysis", self.styles['SectionHeader']))
        
        # Add technical indicators summary
        if 'indicators' in technical:
            story.append(Paragraph("Current Technical Indicators:", self.styles['Heading3']))
            for indicator, value in technical['indicators'].items():
                if isinstance(value, (int, float)):
                    story.append(Paragraph(f"• {indicator}: {value:.2f}", self.styles['Normal']))
                else:
                    story.append(Paragraph(f"• {indicator}: {value}", self.styles['Normal']))
        
        # Add signals summary
        if 'signals' in technical:
            story.append(Spacer(1, 12))
            story.append(Paragraph("Trading Signals:", self.styles['Heading3']))
            story.append(Paragraph(technical['signals'], self.styles['Normal']))
        
        story.append(Spacer(1, 12))
        return story
    
    def _add_charts_section(self, story: List, charts: List[str]) -> List:
        """Add charts section."""
        story.append(Paragraph("Charts and Visualizations", self.styles['SectionHeader']))
        
        for chart_path in charts:
            if os.path.exists(chart_path):
                try:
                    img = Image(chart_path, width=6*inch, height=4*inch)
                    story.append(img)
                    story.append(Spacer(1, 12))
                except Exception as e:
                    story.append(Paragraph(f"Error loading chart: {chart_path}", self.styles['Normal']))
        
        return story
    
    def _add_risk_section(self, story: List, risk_analysis: Dict[str, Any]) -> List:
        """Add risk analysis section."""
        story.append(Paragraph("Risk Analysis", self.styles['SectionHeader']))
        
        risk_data = [
            ['Risk Metric', 'Value'],
            ['Volatility', f"{risk_analysis.get('volatility', 0):.2f}%"],
            ['Value at Risk (95%)', f"{risk_analysis.get('var_95', 0):.2f}%"],
            ['Conditional VaR (95%)', f"{risk_analysis.get('cvar_95', 0):.2f}%"],
            ['Beta', f"{risk_analysis.get('beta', 'N/A')}"],
            ['Tracking Error', f"{risk_analysis.get('tracking_error', 'N/A')}"]
        ]
        
        table = Table(risk_data, colWidths=[2.5*inch, 2.5*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.darkred),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.mistyrose),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        story.append(table)
        story.append(Spacer(1, 12))
        return story
    
    def _add_recommendations_section(self, story: List, recommendations: str) -> List:
        """Add recommendations section."""
        story.append(Paragraph("Investment Recommendations", self.styles['SectionHeader']))
        story.append(Paragraph(recommendations, self.styles['Normal']))
        return story
    
    def get_parameters(self) -> Dict[str, Any]:
        return {
            'report_data': {'type': 'dict', 'required': True, 'description': 'Report data'},
            'report_title': {'type': 'str', 'required': False, 'default': 'Investment Analysis Report',
                           'description': 'Report title'},
            'save_path': {'type': 'str', 'required': False, 'description': 'Path to save the report'}
        }


class PowerPointGenerator(BaseTool):
    """Generate PowerPoint presentations."""
    
    def __init__(self, output_dir: str = "presentations"):
        super().__init__("powerpoint_generator", ToolType.OUTPUT)
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
    
    def execute(self, presentation_data: Dict[str, Any], 
                presentation_title: str = "Investment Analysis", 
                save_path: str = None, **kwargs) -> ToolResult:
        """Generate a PowerPoint presentation."""
        try:
            if save_path is None:
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                save_path = os.path.join(self.output_dir, f"investment_presentation_{timestamp}.pptx")
            
            prs = Presentation()
            
            # Title slide
            self._add_title_slide(prs, presentation_title, 
                                presentation_data.get('subtitle', 'Financial Analysis Report'))
            
            # Executive Summary slide
            if 'executive_summary' in presentation_data:
                self._add_content_slide(prs, "Executive Summary", 
                                      presentation_data['executive_summary'])
            
            # Stock Overview slide
            if 'stock_info' in presentation_data:
                self._add_stock_overview_slide(prs, presentation_data['stock_info'])
            
            # Performance slides
            if 'performance_metrics' in presentation_data:
                self._add_performance_slide(prs, presentation_data['performance_metrics'])
            
            # Technical Analysis slide
            if 'technical_analysis' in presentation_data:
                self._add_technical_slide(prs, presentation_data['technical_analysis'])
            
            # Charts slides
            if 'charts' in presentation_data:
                self._add_chart_slides(prs, presentation_data['charts'])
            
            # Risk Analysis slide
            if 'risk_analysis' in presentation_data:
                self._add_risk_slide(prs, presentation_data['risk_analysis'])
            
            # Recommendations slide
            if 'recommendations' in presentation_data:
                self._add_content_slide(prs, "Investment Recommendations", 
                                      presentation_data['recommendations'])
            
            prs.save(save_path)
            
            return ToolResult(
                success=True,
                data={'presentation_path': save_path},
                metadata={'presentation_type': 'powerpoint', 'title': presentation_title}
            )
        
        except Exception as e:
            return ToolResult(
                success=False,
                error=f"PowerPoint generation failed: {str(e)}"
            )
    
    def _add_title_slide(self, prs, title: str, subtitle: str):
        """Add title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"{subtitle}\n{datetime.now().strftime('%B %d, %Y')}"
    
    def _add_content_slide(self, prs, title: str, content: str):
        """Add a content slide with title and bullet points."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        content_shape.text = content
    
    def _add_stock_overview_slide(self, prs, stock_info: Dict[str, Any]):
        """Add stock overview slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = f"Stock Overview: {stock_info.get('symbol', 'N/A')}"
        
        overview_text = f"""
        • Company: {stock_info.get('company_name', 'N/A')}
        • Sector: {stock_info.get('sector', 'N/A')}
        • Industry: {stock_info.get('industry', 'N/A')}
        • Current Price: ${stock_info.get('current_price', 0):.2f}
        • Market Cap: {stock_info.get('market_cap', 'N/A')}
        """
        
        content_shape.text = overview_text
    
    def _add_performance_slide(self, prs, performance: Dict[str, Any]):
        """Add performance metrics slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Performance Metrics"
        
        perf_text = f"""
        • Total Return: {performance.get('total_return', 0):.2f}%
        • Annualized Return: {performance.get('annualized_return', 0):.2f}%
        • Sharpe Ratio: {performance.get('sharpe_ratio', 0):.3f}
        • Maximum Drawdown: {performance.get('max_drawdown', 0):.2f}%
        • Volatility: {performance.get('volatility', 0):.2f}%
        • Win Rate: {performance.get('win_rate', 0):.2f}%
        """
        
        content_shape.text = perf_text
    
    def _add_technical_slide(self, prs, technical: Dict[str, Any]):
        """Add technical analysis slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Technical Analysis"
        
        # Build technical content
        tech_content = "Key Technical Indicators:\n"
        if 'indicators' in technical:
            for indicator, value in technical['indicators'].items():
                if isinstance(value, (int, float)):
                    tech_content += f"• {indicator}: {value:.2f}\n"
                else:
                    tech_content += f"• {indicator}: {value}\n"
        
        if 'signals' in technical:
            tech_content += f"\nTrading Signals:\n{technical['signals']}"
        
        content_shape.text = tech_content
    
    def _add_chart_slides(self, prs, charts: List[str]):
        """Add chart slides."""
        for chart_path in charts:
            if os.path.exists(chart_path):
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = os.path.basename(chart_path).replace('_', ' ').title()
                
                # Add chart image
                try:
                    slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), 
                                           width=Inches(8), height=Inches(5.5))
                except Exception as e:
                    # If image can't be added, add error text
                    error_box = slide.shapes.add_textbox(Inches(3), Inches(4), Inches(4), Inches(1))
                    error_frame = error_box.text_frame
                    error_frame.text = f"Error loading chart: {str(e)}"
    
    def _add_risk_slide(self, prs, risk_analysis: Dict[str, Any]):
        """Add risk analysis slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Risk Analysis"
        
        risk_text = f"""
        • Volatility: {risk_analysis.get('volatility', 0):.2f}%
        • Value at Risk (95%): {risk_analysis.get('var_95', 0):.2f}%
        • Conditional VaR (95%): {risk_analysis.get('cvar_95', 0):.2f}%
        • Beta: {risk_analysis.get('beta', 'N/A')}
        • Tracking Error: {risk_analysis.get('tracking_error', 'N/A')}
        """
        
        content_shape.text = risk_text
    
    def get_parameters(self) -> Dict[str, Any]:
        return {
            'presentation_data': {'type': 'dict', 'required': True, 'description': 'Presentation data'},
            'presentation_title': {'type': 'str', 'required': False, 'default': 'Investment Analysis',
                                 'description': 'Presentation title'},
            'save_path': {'type': 'str', 'required': False, 'description': 'Path to save the presentation'}
        }